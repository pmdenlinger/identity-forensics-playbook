#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import argparse, os

AUTHOR_NAME = "Paul Denlinger"
AUTHOR_EMAIL = "paul.denlinger@gmail.com"

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--custodian")
    ap.add_argument("--home", help="Home tenant (e.g., contoso.onmicrosoft.com)")
    ap.add_argument("--audience", help="Audience/Resource tenant (e.g., partner.onmicrosoft.com)")
    ap.add_argument("--mfa", help="MFA state/notes")
    ap.add_argument("--ca", help="Conditional Access outcome/notes")
    ap.add_argument("--risky", help="Risky sign-ins count/flag")
    ap.add_argument("--apps", help="Apps used (Teams, SharePoint, Exchange)")
    ap.add_argument("--devices", help="Devices/contexts (corp laptop, iOS, unmanaged)")
    ap.add_argument("--from", dest="dt_from", help="UTC From")
    ap.add_argument("--to", dest="dt_to", help="UTC To")
    ap.add_argument("--notes", help="Free-form notes")
    args = ap.parse_args()

    custodian = args.custodian or input("Custodian: ").strip()
    home = args.home or input("Home tenant: ").strip()
    audience = args.audience or input("Audience/Resource tenant: ").strip()
    mfa = args.mfa or input("MFA: ").strip()
    ca = args.ca or input("Conditional Access: ").strip()
    risky = args.risky or input("Risky sign-ins: ").strip()
    apps = args.apps or input("Apps: ").strip()
    devices = args.devices or input("Devices: ").strip()
    dt_from = args.dt_from or input("UTC From: ").strip()
    dt_to = args.dt_to or input("UTC To: ").strip()
    notes = args.notes or input("Notes: ").strip()

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    blue = RGBColor(0,94,184); light_gray = RGBColor(238,242,246); gray = RGBColor(96,96,96)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    band.fill.solid(); band.fill.fore_color.rgb = light_gray; band.line.color.rgb = blue
    t = band.text_frame; t.clear()
    p = t.paragraphs[0]; r = p.add_run()
    r.text = f"Entra ID Signals — {custodian}"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = blue
    box = slide.shapes.add_textbox(Inches(8.0), Inches(0.35), Inches(4.8), Inches(0.6))
    tf = box.text_frame; tf.clear()
    pn = tf.paragraphs[0]; rn = pn.add_run(); rn.text = f"{AUTHOR_NAME}  •  {AUTHOR_EMAIL}"
    pn.font.size = Pt(14); pn.font.bold = True; pn.font.color.rgb = blue; pn.alignment = PP_ALIGN.RIGHT

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.5), Inches(4.8))
    panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(246,250,255); panel.line.color.rgb = gray
    tf = panel.text_frame; tf.clear()

    rows = [
        f"Home tenant: {home}",
        f"Audience/Resource tenant: {audience}",
        f"MFA: {mfa}",
        f"Conditional Access: {ca}",
        f"Risky sign-ins: {risky}",
        f"Apps: {apps}",
        f"Devices: {devices}",
        f"Window (UTC): {dt_from} → {dt_to}",
        f"Notes: {notes}"
    ]

    for i, row in enumerate(rows):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run(); r.text = row
        p.font.size = Pt(14); p.font.color.rgb = RGBColor(40,40,40)

    outp_dir = "templates"
    os.makedirs(outp_dir, exist_ok=True)
    outp = os.path.join(outp_dir, "entra-id-signals-poster.pptx")
    prs.save(outp)
    print(f"Wrote {outp}")

if __name__ == "__main__":
    main()
