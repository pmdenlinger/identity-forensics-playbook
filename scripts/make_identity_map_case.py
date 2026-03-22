#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import argparse, os

AUTHOR_NAME = "Paul Denlinger"
AUTHOR_EMAIL = "paul.denlinger@gmail.com"

def build_base_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    blue = RGBColor(0,94,184); light_blue = RGBColor(198,224,255)
    gray = RGBColor(96,96,96); light_gray = RGBColor(238,242,246)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    band.fill.solid(); band.fill.fore_color.rgb = light_gray; band.line.color.rgb = blue
    t = band.text_frame; t.clear()
    p = t.paragraphs[0]; r = p.add_run()
    r.text = "Identity Map – Custodian Overview (Work / Historic / Guest / Personal)"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = blue
    box = slide.shapes.add_textbox(Inches(8.0), Inches(0.35), Inches(4.8), Inches(0.6))
    tf = box.text_frame; tf.clear()
    pn = tf.paragraphs[0]; rn = pn.add_run(); rn.text = f"{AUTHOR_NAME}  •  {AUTHOR_EMAIL}"
    pn.font.size = Pt(14); pn.font.bold = True; pn.font.color.rgb = blue; pn.alignment = PP_ALIGN.RIGHT
    for text, left in [("Identities",0.5), ("Evidence Stores",6.5)]:
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(1.2), Inches(5.8), Inches(0.4))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = light_gray; hdr.line.color.rgb = blue
        ht = hdr.text_frame; ht.clear()
        p = ht.paragraphs[0]; r = p.add_run(); r.text = text
        p.font.bold = True; p.font.size = Pt(14)
    legend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.3), Inches(1.2), Inches(3.0), Inches(1.7))
    legend.fill.solid(); legend.fill.fore_color.rgb = light_gray; legend.line.color.rgb = gray
    lt = legend.text_frame; lt.text = "Legend"
    lt.paragraphs[0].font.bold = True; lt.paragraphs[0].font.size = Pt(16)
    items=[(1.65,RGBColor(0,94,184),"Work Identity (Entra ID)"),
           (1.95,RGBColor(0,140,140),"Historic UPN / Alias"),
           (2.25,RGBColor(0,153,0),"Guest (External B2B)"),
           (2.55,RGBColor(153,0,0),"Personal Microsoft Account (MSA)")]
    for y,color,label in items:
        chip=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(y), Inches(0.35), Inches(0.25))
        chip.fill.solid(); chip.fill.fore_color.rgb=color; chip.line.color.rgb=color
        tx=slide.shapes.add_textbox(Inches(10.8), Inches(y-0.02), Inches(2.3), Inches(0.35))
        tf=tx.text_frame; tf.clear()
        p=tf.paragraphs[0]; r=p.add_run(); r.text=label; p.font.size=Pt(12)
    card_specs=[("Work (Current UPN)","name@company.com\nMailbox, Teams, OneDrive",RGBColor(0,94,184),1.7),
                ("Historic UPN / Alias","name@old-company.com\nLegacy mailbox/drive",RGBColor(0,140,140),3.0),
                ("Guest (External B2B)","name_partner.com#EXT#@company.onmicrosoft.com\nExternal Teams/Sites",RGBColor(0,153,0),4.3),
                ("Personal MSA (If Signaled)","name@outlook.com\nLink access / mobile",RGBColor(153,0,0),5.6)]
    cards=[]
    for title,body,color,y in card_specs:
        rect=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(5.8), Inches(1.0))
        rect.fill.solid(); rect.fill.fore_color.rgb=light_blue; rect.line.color.rgb=color
        tf=rect.text_frame; tf.clear()
        p1=tf.paragraphs[0]; r1=p1.add_run(); r1.text=title
        p1.font.bold=True; p1.font.size=Pt(16); p1.font.color.rgb=color
        p2=tf.add_paragraph(); p2.text=body; p2.font.size=Pt(12)
        cards.append((rect,color))
    stores=[]; gray = RGBColor(96,96,96)
    s_specs=[("Exchange Online","Primary + historic mailboxes\nForwarding/redirect rules",1.7),
             ("Teams","Internal channels + Partner (Guest) channels\nMessage exports / transcripts",3.0),
             ("SharePoint / OneDrive","Project sites + personal drives\nShare-link types & external recipients",4.3),
             ("Audit & OAuth","Entra sign-ins / CA results\nOAuth grants / device signals",5.6)]
    for title,body,y in s_specs:
        rect=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(y), Inches(5.8), Inches(1.0))
        rect.fill.solid(); rect.fill.fore_color.rgb=RGBColor(246,250,255); rect.line.color.rgb=gray
        tf=rect.text_frame; tf.clear()
        p1=tf.paragraphs[0]; r1=p1.add_run(); r1.text=title
        p1.font.bold=True; p1.font.size=Pt(16); p1.font.color.rgb=gray
        p2=tf.add_paragraph(); p2.text=body; p2.font.size=Pt(12)
        stores.append(rect)
    for i,(rect,color) in enumerate(cards):
        sx=rect.left+rect.width; sy=rect.top+rect.height/2
        tgt=stores[i if i<len(stores) else -1]
        tx=tgt.left; ty=tgt.top+tgt.height/2
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx,sy,tx,ty)
        line.line.color.rgb=color
    ft=slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.5)).text_frame
    p=ft.paragraphs[0]; r=p.add_run()
    r.text=("Identity Forensics maps Work / Historic / Guest / Personal identities to the evidence stores they touched. "
            "Scope expands ONLY where signals justify it (sign-ins, memberships, share links, rules, OAuth). "
            "Time-bounded. Hashed. Manifested.")
    p.font.size=Pt(11); p.font.color.rgb=gray
    return slide, blue, light_gray, gray

def add_case_panel(slide, blue, light_gray, case_id, custodian, dt_from, dt_to):
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(0.3), Inches(3.6), Inches(0.8))
    panel.fill.solid(); panel.fill.fore_color.rgb = light_gray; panel.line.color.rgb = blue
    tf = panel.text_frame; tf.clear()
    rows = [f"Case: {case_id or ''}",
            f"Custodian: {custodian or ''}",
            f"UTC From: {dt_from or ''}",
            f"UTC To: {dt_to or ''}"]
    for i, row in enumerate(rows):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run(); r.text = row
        p.font.size = Pt(12); p.font.color.rgb = RGBColor(40,40,40)

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--case", dest="case_id")
    ap.add_argument("--custodian", dest="custodian")
    ap.add_argument("--from", dest="dt_from")
    ap.add_argument("--to", dest="dt_to")
    args = ap.parse_args()
    case_id = args.case_id or input("Case ID: ").strip()
    custodian = args.custodian or input("Custodian: ").strip()
    dt_from = args.dt_from or input("UTC From (e.g., 2026-01-01T00:00Z): ").strip()
    dt_to = args.dt_to or input("UTC To   (e.g., 2026-02-29T23:59Z): ").strip()
    os.makedirs("templates", exist_ok=True)
    prs = Presentation()
    slide, blue, light_gray, gray = build_base_slide(prs)
    add_case_panel(slide, blue, light_gray, case_id, custodian, dt_from, dt_to)
    outp = "templates/identity-map-poster_case.pptx"
    prs.save(outp)
    print(f"Wrote {outp}")

if __name__=="__main__":
    main()
